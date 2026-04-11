"""Document export endpoints.

Generates downloadable DOCX, PPTX, and printable HTML documents
from brief/signal content passed in the request body.

Performance guardrails:
  - Asyncio semaphore limits concurrent in-flight exports.
  - Per-user rate limiting via slowapi (10 exports / minute).
  - Section count and content length caps prevent memory exhaustion.
"""

import asyncio
import html as _html
import io
import logging
from datetime import datetime
from typing import Literal

from fastapi import APIRouter, Depends, HTTPException, Request, status
from fastapi.responses import HTMLResponse, Response
from pydantic import BaseModel, Field

from sqlalchemy.ext.asyncio import AsyncSession

from backend.auth.dependencies import get_current_user
from backend.auth.rate_limit import limiter
from backend.auth.schemas import AuthContext
from backend.database import get_db
from backend.middleware.feature_gating import get_current_organization
from backend.models.organization import Organization
from backend.services.credit_service import CreditService, InsufficientCreditsError

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/exports")

# ── Limits ─────────────────────────────────────────────────────────────────────

MAX_EXPORT_SECTIONS = 50
MAX_SECTION_CONTENT_LENGTH = 50_000  # characters per section
MAX_TOTAL_CONTENT_LENGTH = 500_000  # characters total across all sections
MAX_CONCURRENT_EXPORTS = 5  # at most 5 exports generating simultaneously

_export_semaphore = asyncio.Semaphore(MAX_CONCURRENT_EXPORTS)

# ── Request schema ─────────────────────────────────────────────────────────────


class ContentSection(BaseModel):
    heading: str = Field(..., max_length=500)
    content: str = Field(..., max_length=MAX_SECTION_CONTENT_LENGTH)


class BriefExportRequest(BaseModel):
    title: str = Field(..., max_length=500)
    subtitle: str | None = Field(None, max_length=500)
    domain: str | None = Field(None, max_length=200)
    author: str | None = Field(None, max_length=200)
    confidence: int | None = Field(None, ge=0, le=100)
    sections: list[ContentSection] = Field(default=[], max_length=MAX_EXPORT_SECTIONS)
    format: Literal["docx", "pptx", "pdf-html"] = "docx"


# ── Helpers ────────────────────────────────────────────────────────────────────


def _safe_filename(title: str) -> str:
    """Create a safe filename from a title string."""
    return (
        title.lower()
        .replace(" ", "-")
        .replace("/", "-")
        .replace("\\", "-")
        .replace(":", "")
        .replace('"', "")
        .replace("'", "")[:60]
    )


def _build_docx(req: BriefExportRequest) -> bytes:
    """Generate a .docx document from brief content."""
    from docx import Document  # type: ignore[import]
    from docx.shared import Pt, RGBColor  # type: ignore[import]

    doc = Document()

    # Title
    title_para = doc.add_heading(req.title, level=0)
    title_para.style.font.size = Pt(20)

    # Subtitle
    if req.subtitle:
        sub = doc.add_paragraph(req.subtitle)
        sub.style.font.size = Pt(12)
        sub.style.font.italic = True

    # Metadata line
    meta_parts = []
    if req.domain:
        meta_parts.append(f"Domain: {req.domain}")
    if req.author:
        meta_parts.append(f"Author: {req.author}")
    if req.confidence is not None:
        meta_parts.append(f"Confidence: {req.confidence}%")
    meta_parts.append(f"Generated: {datetime.utcnow().strftime('%d %B %Y')}")
    meta_para = doc.add_paragraph(" · ".join(meta_parts))
    meta_para.style.font.size = Pt(10)
    meta_para.style.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    doc.add_paragraph()  # spacer

    # Sections
    for section in req.sections:
        doc.add_heading(section.heading, level=1)
        doc.add_paragraph(section.content)

    # Footer
    doc.add_paragraph()
    footer = doc.add_paragraph(
        f"Cogent Intelligence — Proprietary Report · {datetime.utcnow().strftime('%d %B %Y')}"
    )
    footer.style.font.size = Pt(9)
    footer.style.font.italic = True

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_pptx(req: BriefExportRequest) -> bytes:
    """Generate a .pptx presentation from brief content."""
    from pptx import Presentation  # type: ignore[import]
    from pptx.dml.color import RGBColor  # type: ignore[import]
    from pptx.enum.text import PP_ALIGN  # type: ignore[import]
    from pptx.util import Inches, Pt  # type: ignore[import]

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # fully blank

    # ── Title slide ──────────────────────────────────────────────────────────
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = req.title
    if req.subtitle:
        title_slide.placeholders[1].text = req.subtitle

    # ── Meta slide ───────────────────────────────────────────────────────────
    meta_slide = prs.slides.add_slide(blank_layout)
    meta_box = meta_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12), Inches(6)
    )
    tf = meta_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Report Details"
    p.font.size = Pt(20)
    p.font.bold = True

    meta_parts = []
    if req.domain:
        meta_parts.append(f"Domain: {req.domain}")
    if req.author:
        meta_parts.append(f"Author: {req.author}")
    if req.confidence is not None:
        meta_parts.append(f"Confidence: {req.confidence}%")
    meta_parts.append(f"Date: {datetime.utcnow().strftime('%d %B %Y')}")

    for part in meta_parts:
        p2 = tf.add_paragraph()
        p2.text = part
        p2.font.size = Pt(14)

    # ── One slide per section ─────────────────────────────────────────────────
    for section in req.sections:
        slide = prs.slides.add_slide(blank_layout)

        # Section heading
        heading_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
        )
        heading_tf = heading_box.text_frame
        heading_para = heading_tf.paragraphs[0]
        heading_para.text = section.heading
        heading_para.font.size = Pt(22)
        heading_para.font.bold = True
        heading_para.font.color.rgb = RGBColor(0x10, 0x18, 0x36)

        # Section body
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(12), Inches(5.8)
        )
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_para = body_tf.paragraphs[0]
        body_para.text = section.content
        body_para.font.size = Pt(14)

    # ── Final slide ───────────────────────────────────────────────────────────
    final_slide = prs.slides.add_slide(blank_layout)
    footer_box = final_slide.shapes.add_textbox(
        Inches(1), Inches(3), Inches(11), Inches(1.5)
    )
    footer_tf = footer_box.text_frame
    footer_para = footer_tf.paragraphs[0]
    footer_para.text = "Cogent Intelligence"
    footer_para.font.size = Pt(28)
    footer_para.font.bold = True
    footer_para.alignment = PP_ALIGN.CENTER

    p2 = footer_tf.add_paragraph()
    p2.text = f"Proprietary Report · {datetime.utcnow().strftime('%d %B %Y')}"
    p2.font.size = Pt(14)
    p2.alignment = PP_ALIGN.CENTER

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_pdf_html(req: BriefExportRequest) -> str:
    """Generate a printable HTML page for browser print-to-PDF."""
    e = _html.escape  # shorthand for readability
    sections_html = "".join(
        f"<h2>{e(s.heading)}</h2><p>{e(s.content)}</p>" for s in req.sections
    )
    meta_parts = []
    if req.domain:
        meta_parts.append(f"<strong>Domain:</strong> {e(req.domain)}")
    if req.author:
        meta_parts.append(f"<strong>Author:</strong> {e(req.author)}")
    if req.confidence is not None:
        meta_parts.append(f"<strong>Confidence:</strong> {req.confidence}%")
    meta_parts.append(
        f"<strong>Generated:</strong> {datetime.utcnow().strftime('%d %B %Y')}"
    )
    meta_html = " &nbsp;·&nbsp; ".join(meta_parts)

    escaped_title = e(req.title)
    escaped_subtitle = e(req.subtitle) if req.subtitle else None

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <title>{escaped_title}</title>
  <style>
    *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
    body {{
      font-family: 'Georgia', serif;
      color: #1a1a2e;
      background: #fff;
      padding: 48px 64px;
      max-width: 860px;
      margin: 0 auto;
    }}
    h1 {{ font-size: 28px; font-weight: 700; margin-bottom: 8px; line-height: 1.3; }}
    .subtitle {{ font-size: 16px; color: #555; font-style: italic; margin-bottom: 12px; }}
    .meta {{ font-size: 12px; color: #777; margin-bottom: 32px; border-bottom: 1px solid #e5e7eb; padding-bottom: 12px; }}
    h2 {{ font-size: 17px; font-weight: 600; margin: 28px 0 8px; color: #1a1a2e; }}
    p {{ font-size: 14px; line-height: 1.75; color: #333; margin-bottom: 12px; }}
    .footer {{ margin-top: 48px; padding-top: 16px; border-top: 1px solid #e5e7eb; font-size: 11px; color: #aaa; font-style: italic; }}
    @media print {{
      body {{ padding: 24px 32px; }}
      h1 {{ page-break-after: avoid; }}
      h2 {{ page-break-after: avoid; }}
    }}
  </style>
</head>
<body>
  <h1>{escaped_title}</h1>
  {"<p class='subtitle'>" + escaped_subtitle + "</p>" if escaped_subtitle else ""}
  <p class="meta">{meta_html}</p>
  {sections_html}
  <div class="footer">Cogent Intelligence — Proprietary Report</div>
  <script>window.onload = function() {{ window.print(); }}</script>
</body>
</html>"""


# ── Endpoint ───────────────────────────────────────────────────────────────────


@router.post("/brief")
@limiter.limit("10/minute")
async def export_brief(
    request: Request,
    body: BriefExportRequest,
    auth: AuthContext = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
    organization: Organization = Depends(get_current_organization),
) -> Response:
    """Export a brief or signal document as DOCX, PPTX, or printable HTML (PDF).

    - **docx**: Returns a Word document (.docx)
    - **pptx**: Returns a PowerPoint presentation (.pptx)
    - **pdf-html**: Returns a print-ready HTML page (open in browser, Ctrl+P to save as PDF)

    Rate limited to 10 exports per minute per user.
    Concurrency limited to 5 simultaneous exports across all users.
    """
    credit_service = CreditService(db)
    try:
        await credit_service.consume_credits(
            org_id=organization.id,
            user_id=auth.user_id,
            action_type="document_export",
            credits=5,
            metadata={"format": body.format, "section_count": len(body.sections)},
        )
    except InsufficientCreditsError as e:
        raise HTTPException(
            status_code=402,
            detail=(
                f"Insufficient credits for export. "
                f"Requires {e.required} credits and {e.remaining} remain."
            ),
        ) from e

    # Validate total content size
    total_chars = sum(len(s.content) for s in body.sections)
    if total_chars > MAX_TOTAL_CONTENT_LENGTH:
        raise HTTPException(
            status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
            detail=(
                f"Total content length ({total_chars:,} chars) exceeds "
                f"maximum ({MAX_TOTAL_CONTENT_LENGTH:,} chars)"
            ),
        )

    # Acquire concurrency semaphore (fail fast if too many exports in-flight)
    if _export_semaphore.locked():
        raise HTTPException(
            status_code=status.HTTP_429_TOO_MANY_REQUESTS,
            detail="Export service is busy. Please try again in a few seconds.",
        )

    async with _export_semaphore:
        safe_name = _safe_filename(body.title)

        if body.format == "docx":
            content = _build_docx(body)
            return Response(
                content=content,
                media_type=(
                    "application/vnd.openxmlformats-officedocument"
                    ".wordprocessingml.document"
                ),
                headers={
                    "Content-Disposition": f'attachment; filename="{safe_name}.docx"'
                },
            )

        if body.format == "pptx":
            content = _build_pptx(body)
            return Response(
                content=content,
                media_type=(
                    "application/vnd.openxmlformats-officedocument"
                    ".presentationml.presentation"
                ),
                headers={
                    "Content-Disposition": f'attachment; filename="{safe_name}.pptx"'
                },
            )

        # pdf-html — return printable HTML that triggers window.print()
        html = _build_pdf_html(body)
        return HTMLResponse(
            content=html,
            headers={
                "Content-Disposition": f'inline; filename="{safe_name}-print.html"'
            },
        )
