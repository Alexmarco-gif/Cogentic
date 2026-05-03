"""Background job handlers used by RQ workers and scheduler hooks.

This module contains the synchronous entry points enqueued by ``backend.job_queue``.
Each public function wraps an async implementation so worker processes can execute
real document analysis, batch import, analytics reporting, webhook delivery, and
notification delivery without returning placeholder success payloads.
"""

from __future__ import annotations

import asyncio
import csv
import hashlib
import hmac
import io
import json
import logging
import mimetypes
import re
import socket
import zipfile
from collections import Counter
from datetime import UTC, datetime, timedelta
from html import unescape
from ipaddress import ip_address
from pathlib import Path
from typing import Any
from urllib.parse import urlparse
from uuid import UUID
from xml.etree import ElementTree

import httpx
from selectolax.parser import HTMLParser
from sqlalchemy import delete, func, select

from backend.database import AsyncSessionLocal
from backend.models.ai_job import AIJob
from backend.models.credit_transaction import CreditTransaction
from backend.models.document import Document
from backend.models.intelligence_brief import IntelligenceBrief
from backend.models.org_user import OrgUser
from backend.models.organization import Organization
from backend.models.signal import Signal
from backend.services.email_service import (
    send_data_export_request_email,
    send_deletion_request_email,
    send_email,
)
from backend.services.feedback_retraining import run_feedback_retraining_job
from backend.storage import (
    delete_storage_object,
    download_storage_object,
    local_path_from_storage_path,
    storage_path_is_object,
)

logger = logging.getLogger(__name__)

_BLOCKED_WEBHOOK_HOSTS = {
    "localhost",
    "127.0.0.1",
    "0.0.0.0",
    "::1",
    "[::1]",
    "metadata.google.internal",
}
_TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".html",
    ".htm",
    ".xml",
    ".yaml",
    ".yml",
    ".log",
    ".rtf",
}
_STOPWORDS = {
    "a",
    "an",
    "and",
    "are",
    "as",
    "at",
    "be",
    "by",
    "for",
    "from",
    "has",
    "in",
    "is",
    "it",
    "its",
    "of",
    "on",
    "or",
    "that",
    "the",
    "their",
    "this",
    "to",
    "was",
    "were",
    "with",
}
_DATE_PATTERNS = (
    re.compile(r"\b\d{4}-\d{2}-\d{2}\b"),
    re.compile(r"\b\d{1,2}/\d{1,2}/\d{2,4}\b"),
)
_EMAIL_RE = re.compile(r"\b[A-Z0-9._%+-]+@[A-Z0-9.-]+\.[A-Z]{2,}\b", re.IGNORECASE)
_URL_RE = re.compile(r"https?://[^\s<>\"]+", re.IGNORECASE)
_TOKEN_RE = re.compile(r"[A-Za-z][A-Za-z0-9_-]{2,}")
_MAX_DOCUMENT_BYTES = 10 * 1024 * 1024
_MAX_EXTRACTED_TEXT = 200_000
_MAX_REMOTE_DOWNLOAD_BYTES = 10 * 1024 * 1024


def _safe_uuid(value: str | UUID) -> UUID:
    return value if isinstance(value, UUID) else UUID(str(value))


def _utcnow() -> datetime:
    return datetime.now(UTC)


def _coerce_text(raw: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _truncate_text(text: str, limit: int = _MAX_EXTRACTED_TEXT) -> str:
    cleaned = re.sub(r"\s+", " ", text).strip()
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: limit - 3].rstrip() + "..."


def _html_to_text(text: str) -> str:
    parser = HTMLParser(text)
    if parser.body is not None:
        extracted = parser.body.text(separator=" ", strip=True)
    else:
        extracted = parser.text(separator=" ", strip=True)
    return unescape(extracted)


def _xml_to_text(text: str) -> str:
    try:
        root = ElementTree.fromstring(text)
        return " ".join(part.strip() for part in root.itertext() if part.strip())
    except ElementTree.ParseError:
        return text


def _extract_docx_text(raw: bytes) -> str:
    try:
        from docx import Document as DocxDocument  # type: ignore[import]
    except ImportError as exc:  # pragma: no cover - dependency should exist
        raise RuntimeError("DOCX extraction requires python-docx") from exc

    doc = DocxDocument(io.BytesIO(raw))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx_text(raw: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore[import]
    except ImportError as exc:  # pragma: no cover - dependency should exist
        raise RuntimeError("PPTX extraction requires python-pptx") from exc

    presentation = Presentation(io.BytesIO(raw))
    slides: list[str] = []
    for slide in presentation.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                texts.append(text.strip())
        if texts:
            slides.append("\n".join(texts))
    return "\n\n".join(slides)


def _extract_pdf_text(raw: bytes) -> str:
    try:
        from pypdf import PdfReader  # type: ignore[import]
    except ImportError as exc:
        raise RuntimeError(
            "PDF extraction requires the optional pypdf dependency"
        ) from exc

    reader = PdfReader(io.BytesIO(raw))
    pages = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return "\n".join(page.strip() for page in pages if page.strip())


def _extract_text_from_bytes(
    raw: bytes,
    *,
    filename: str,
    content_type: str | None,
) -> str:
    extension = Path(filename).suffix.lower()
    detected = (content_type or "").split(";")[0].strip().lower()

    if len(raw) > _MAX_DOCUMENT_BYTES:
        raise RuntimeError(
            f"Document exceeds maximum supported size of {_MAX_DOCUMENT_BYTES} bytes"
        )

    if extension == ".docx" or detected.endswith(
        "vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        return _extract_docx_text(raw)

    if extension == ".pptx" or detected.endswith(
        "vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        return _extract_pptx_text(raw)

    if extension == ".pdf" or detected == "application/pdf":
        return _extract_pdf_text(raw)

    text = _coerce_text(raw)

    if extension in {".html", ".htm"} or detected in {
        "text/html",
        "application/xhtml+xml",
    }:
        return _html_to_text(text)

    if extension == ".xml" or detected in {"text/xml", "application/xml"}:
        return _xml_to_text(text)

    if extension == ".json" or detected == "application/json":
        try:
            parsed = json.loads(text)
            return json.dumps(parsed, indent=2, ensure_ascii=True)
        except json.JSONDecodeError:
            return text

    if extension == ".csv" or detected == "text/csv":
        reader = csv.reader(io.StringIO(text))
        rows = [" | ".join(cell.strip() for cell in row) for row in reader]
        return "\n".join(row for row in rows if row.strip())

    if extension in _TEXT_EXTENSIONS or detected.startswith("text/"):
        return text

    if extension == ".zip":
        with zipfile.ZipFile(io.BytesIO(raw)) as archive:
            names = archive.namelist()
        return "ZIP archive containing:\n" + "\n".join(names)

    raise RuntimeError(
        f"Unsupported document format for analysis: {extension or detected or 'unknown'}"
    )


def _extract_keywords(text: str, limit: int = 10) -> list[str]:
    counts = Counter(
        token.lower()
        for token in _TOKEN_RE.findall(text)
        if token.lower() not in _STOPWORDS and len(token) > 2
    )
    return [word for word, _count in counts.most_common(limit)]


def _build_summary(text: str) -> str:
    sentences = re.split(r"(?<=[.!?])\s+", text)
    summary = " ".join(part.strip() for part in sentences[:3] if part.strip())
    if not summary:
        summary = text[:320]
    return _truncate_text(summary, limit=500)


def _classify_document(
    filename: str, content_type: str | None, text: str
) -> dict[str, Any]:
    extension = Path(filename).suffix.lower().lstrip(".") or "unknown"
    lowered = text.lower()
    tags: list[str] = []
    if "invoice" in lowered:
        tags.append("invoice")
    if "contract" in lowered or "agreement" in lowered:
        tags.append("contract")
    if "market" in lowered or "competitor" in lowered:
        tags.append("market_intelligence")
    if "financial" in lowered or "revenue" in lowered:
        tags.append("finance")
    return {
        "file_extension": extension,
        "content_type": content_type
        or mimetypes.guess_type(filename)[0]
        or "application/octet-stream",
        "tags": tags,
        "contains_pii": bool(_EMAIL_RE.search(text)),
    }


def _build_document_analysis(
    *,
    document: Document,
    extracted_text: str,
    content_type: str | None,
    source: str,
) -> dict[str, Any]:
    dates = []
    for pattern in _DATE_PATTERNS:
        dates.extend(match.group(0) for match in pattern.finditer(extracted_text))

    return {
        "summary": _build_summary(extracted_text),
        "word_count": len(re.findall(r"\b\w+\b", extracted_text)),
        "character_count": len(extracted_text),
        "keywords": _extract_keywords(extracted_text),
        "emails": sorted(set(_EMAIL_RE.findall(extracted_text)))[:10],
        "urls": sorted(set(_URL_RE.findall(extracted_text)))[:10],
        "dates": sorted(set(dates))[:10],
        "classification": _classify_document(
            document.filename, content_type, extracted_text
        ),
        "source": source,
        "analyzed_at": _utcnow().isoformat(),
    }


def _is_private_host(hostname: str) -> bool:
    try:
        infos = socket.getaddrinfo(hostname, None)
    except socket.gaierror:
        raise ValueError(f"Unable to resolve webhook hostname: {hostname}")

    for info in infos:
        candidate = info[4][0]
        addr = ip_address(candidate)
        if (
            addr.is_private
            or addr.is_loopback
            or addr.is_link_local
            or addr.is_multicast
            or addr.is_reserved
            or addr.is_unspecified
        ):
            return True
    return False


def _validate_webhook_url(url: str) -> bool:
    try:
        parsed = urlparse(url)
    except Exception:
        return False

    if parsed.scheme not in {"http", "https"}:
        return False

    hostname = (parsed.hostname or "").lower()
    if not hostname or hostname in _BLOCKED_WEBHOOK_HOSTS:
        return False

    try:
        return not _is_private_host(hostname)
    except ValueError:
        return False


async def _read_remote_document(storage_path: str) -> bytes:
    parsed = urlparse(storage_path)
    if parsed.scheme not in {"http", "https"}:
        raise RuntimeError("Only http and https remote document sources are supported")
    if not _validate_webhook_url(storage_path):
        raise RuntimeError("Remote document source is blocked for safety reasons")

    async with httpx.AsyncClient(timeout=20.0, follow_redirects=True) as client:
        response = await client.get(storage_path)
        response.raise_for_status()
        content_length = int(response.headers.get("content-length", "0") or 0)
        if content_length and content_length > _MAX_REMOTE_DOWNLOAD_BYTES:
            raise RuntimeError("Remote document exceeds supported download limit")
        if len(response.content) > _MAX_REMOTE_DOWNLOAD_BYTES:
            raise RuntimeError("Remote document exceeds supported download limit")
        return response.content


async def _read_document_source(storage_path: str) -> tuple[bytes, str]:
    parsed = urlparse(storage_path)

    if storage_path_is_object(storage_path):
        return await asyncio.to_thread(download_storage_object, storage_path), "s3"

    if parsed.scheme in {"http", "https"}:
        return await _read_remote_document(storage_path), "remote"

    path = local_path_from_storage_path(storage_path)
    if not path.exists():
        raise FileNotFoundError(f"Document source not found: {path}")

    return path.read_bytes(), "local"


def _delete_storage_path(storage_path: str | None) -> bool:
    if not storage_path:
        return False

    if storage_path_is_object(storage_path):
        return delete_storage_object(storage_path)

    path = local_path_from_storage_path(storage_path)
    if path.exists():
        path.unlink()
        return True
    return False


async def _set_job_failed(
    db: Any,
    job: AIJob | None,
    document: Document | None,
    *,
    error_message: str,
    started_at: datetime | None = None,
) -> dict[str, Any]:
    if document is not None:
        document.processing_status = "failed"
    if job is not None:
        started = started_at or job.started_at or _utcnow()
        completed = _utcnow()
        job.status = "failed"
        job.error_message = error_message
        job.result = {"status": "failed", "error": error_message}
        job.completed_at = completed
        job.duration_ms = int((completed - started).total_seconds() * 1000)
    await db.commit()
    return {"status": "failed", "error": error_message}


async def _process_document_analysis_async(job_id: str | UUID) -> dict[str, Any]:
    job_uuid = _safe_uuid(job_id)

    async with AsyncSessionLocal() as db:
        result = await db.execute(select(AIJob).where(AIJob.id == job_uuid))
        job = result.scalar_one_or_none()
        if not job:
            raise ValueError(f"AI job not found: {job_uuid}")

        if not job.document_id:
            return await _set_job_failed(
                db, job, None, error_message="AI job is not linked to a document"
            )

        document = await db.get(Document, job.document_id)
        if not document:
            return await _set_job_failed(
                db, job, None, error_message=f"Document not found: {job.document_id}"
            )
        if not document.storage_path:
            return await _set_job_failed(
                db, job, document, error_message="Document storage_path is missing"
            )

        started_at = _utcnow()
        job.status = "running"
        job.started_at = started_at
        job.attempts = (job.attempts or 0) + 1
        document.processing_status = "processing"
        await db.commit()

        try:
            raw, source = await _read_document_source(document.storage_path)
            extracted_text = _truncate_text(
                _extract_text_from_bytes(
                    raw,
                    filename=document.filename,
                    content_type=document.content_type,
                )
            )
            if not extracted_text:
                raise RuntimeError(
                    "No textual content could be extracted from the document"
                )

            analysis = _build_document_analysis(
                document=document,
                extracted_text=extracted_text,
                content_type=document.content_type,
                source=source,
            )

            document.extracted_text = extracted_text
            document.processing_status = "completed"
            completed_at = _utcnow()
            job.status = "completed"
            job.error_message = None
            job.result = {
                "status": "completed",
                "document_id": str(document.id),
                "analysis": analysis,
            }
            job.completed_at = completed_at
            job.duration_ms = int((completed_at - started_at).total_seconds() * 1000)
            await db.commit()
            return job.result
        except Exception as exc:
            logger.exception("Document analysis failed for job %s", job_uuid)
            return await _set_job_failed(
                db,
                job,
                document,
                error_message=str(exc),
                started_at=started_at,
            )


async def _batch_document_import_async(
    org_id: str | UUID,
    user_id: str | UUID,
    documents: list[dict[str, Any] | str],
) -> dict[str, Any]:
    org_uuid = _safe_uuid(org_id)
    user_uuid = _safe_uuid(user_id)

    async with AsyncSessionLocal() as db:
        organization = await db.get(Organization, org_uuid)
        if not organization:
            raise ValueError(f"Organization not found: {org_uuid}")

        membership = await db.execute(
            select(OrgUser).where(
                OrgUser.org_id == org_uuid, OrgUser.user_id == user_uuid
            )
        )
        if membership.scalar_one_or_none() is None:
            raise ValueError("User is not a member of the target organization")

        imported: list[dict[str, Any]] = []
        failed: list[dict[str, Any]] = []

        for entry in documents:
            payload = {"storage_path": entry} if isinstance(entry, str) else dict(entry)
            storage_path = str(payload.get("storage_path") or "").strip()
            filename = str(payload.get("filename") or "").strip()

            if not storage_path:
                failed.append({"status": "failed", "error": "storage_path is required"})
                continue

            try:
                raw, _source = await _read_document_source(storage_path)
                if not filename:
                    parsed = urlparse(storage_path)
                    filename = Path(parsed.path or storage_path).name or (
                        f"document-{len(imported) + len(failed) + 1}"
                    )
                content_type = (
                    payload.get("content_type")
                    or mimetypes.guess_type(filename)[0]
                    or "application/octet-stream"
                )
                size_bytes = int(payload.get("size_bytes") or len(raw))

                document = Document(
                    org_id=org_uuid,
                    owner_id=user_uuid,
                    filename=filename,
                    content_type=content_type,
                    size_bytes=size_bytes,
                    storage_path=storage_path,
                    processing_status="pending",
                )
                db.add(document)
                await db.flush()

                job = AIJob(
                    org_id=org_uuid,
                    user_id=user_uuid,
                    document_id=document.id,
                    job_type="analyze_document",
                    input_params={"storage_path": storage_path, "filename": filename},
                    status="queued",
                )
                db.add(job)
                await db.commit()
                await db.refresh(document)
                await db.refresh(job)

                analysis_result = await _process_document_analysis_async(job.id)
                if analysis_result["status"] == "completed":
                    imported.append(
                        {
                            "status": "success",
                            "document_id": str(document.id),
                            "job_id": str(job.id),
                            "filename": filename,
                        }
                    )
                else:
                    failed.append(
                        {
                            "status": "failed",
                            "document_id": str(document.id),
                            "job_id": str(job.id),
                            "filename": filename,
                            "error": analysis_result.get("error"),
                        }
                    )
            except Exception as exc:
                logger.exception("Batch document import failed for %s", storage_path)
                failed.append(
                    {
                        "status": "failed",
                        "filename": filename or storage_path,
                        "error": str(exc),
                    }
                )

        return {
            "status": "completed",
            "org_id": str(org_uuid),
            "user_id": str(user_uuid),
            "imported_count": len(imported),
            "failed_count": len(failed),
            "imported": imported,
            "failed": failed,
        }


async def _generate_analytics_report_async(
    org_id: str | UUID,
    *,
    report_type: str = "usage",
    days: int = 30,
) -> dict[str, Any]:
    org_uuid = _safe_uuid(org_id)
    window_end = _utcnow()
    window_start = window_end - timedelta(days=max(days, 1))

    async with AsyncSessionLocal() as db:
        organization = await db.get(Organization, org_uuid)
        if not organization:
            raise ValueError(f"Organization not found: {org_uuid}")

        total_documents = await db.scalar(
            select(func.count())
            .select_from(Document)
            .where(Document.org_id == org_uuid, Document.deleted_at.is_(None))
        )
        completed_documents = await db.scalar(
            select(func.count())
            .select_from(Document)
            .where(
                Document.org_id == org_uuid,
                Document.deleted_at.is_(None),
                Document.processing_status == "completed",
            )
        )
        failed_documents = await db.scalar(
            select(func.count())
            .select_from(Document)
            .where(
                Document.org_id == org_uuid,
                Document.deleted_at.is_(None),
                Document.processing_status == "failed",
            )
        )
        total_jobs = await db.scalar(
            select(func.count()).select_from(AIJob).where(AIJob.org_id == org_uuid)
        )
        running_jobs = await db.scalar(
            select(func.count())
            .select_from(AIJob)
            .where(AIJob.org_id == org_uuid, AIJob.status == "running")
        )
        failed_jobs = await db.scalar(
            select(func.count())
            .select_from(AIJob)
            .where(AIJob.org_id == org_uuid, AIJob.status == "failed")
        )
        total_briefs = await db.scalar(
            select(func.count())
            .select_from(IntelligenceBrief)
            .where(
                IntelligenceBrief.org_id == org_uuid,
                IntelligenceBrief.deleted_at.is_(None),
            )
        )
        published_briefs = await db.scalar(
            select(func.count())
            .select_from(IntelligenceBrief)
            .where(
                IntelligenceBrief.org_id == org_uuid,
                IntelligenceBrief.deleted_at.is_(None),
                IntelligenceBrief.status == "published",
            )
        )
        org_signals = await db.scalar(
            select(func.count()).select_from(Signal).where(Signal.org_id == org_uuid)
        )
        recent_signals = await db.scalar(
            select(func.count())
            .select_from(Signal)
            .where(Signal.org_id == org_uuid, Signal.created_at >= window_start)
        )
        credit_rows = await db.execute(
            select(
                CreditTransaction.action_type,
                func.sum(CreditTransaction.credits_consumed),
            )
            .where(
                CreditTransaction.org_id == org_uuid,
                CreditTransaction.created_at >= window_start,
            )
            .group_by(CreditTransaction.action_type)
            .order_by(func.sum(CreditTransaction.credits_consumed).desc())
        )
        credits_by_action = {
            action: int(total or 0) for action, total in credit_rows.all()
        }

        total_credits_used = sum(credits_by_action.values())
        return {
            "status": "completed",
            "report_type": report_type,
            "generated_at": window_end.isoformat(),
            "window": {
                "days": days,
                "start": window_start.isoformat(),
                "end": window_end.isoformat(),
            },
            "organization": {
                "id": str(organization.id),
                "name": organization.name,
                "slug": organization.slug,
                "pricing_tier": organization.pricing_tier,
                "trial_status": organization.trial_status,
            },
            "summary": {
                "documents_total": int(total_documents or 0),
                "documents_completed": int(completed_documents or 0),
                "documents_failed": int(failed_documents or 0),
                "ai_jobs_total": int(total_jobs or 0),
                "ai_jobs_running": int(running_jobs or 0),
                "ai_jobs_failed": int(failed_jobs or 0),
                "briefs_total": int(total_briefs or 0),
                "briefs_published": int(published_briefs or 0),
                "signals_total": int(org_signals or 0),
                "signals_recent": int(recent_signals or 0),
                "credits_used_in_window": int(total_credits_used),
                "credits_allocated_monthly": int(
                    organization.credits_allocated_monthly
                ),
                "credits_consumed": int(organization.credits_consumed),
            },
            "breakdowns": {
                "credits_by_action": credits_by_action,
            },
        }


async def _cleanup_expired_documents_async() -> dict[str, Any]:
    now = _utcnow()
    soft_delete_cutoff = now - timedelta(days=30)

    async with AsyncSessionLocal() as db:
        rows = await db.execute(
            select(Document).where(
                (
                    Document.retention_date.is_not(None)
                    & (Document.retention_date <= now)
                )
                | (
                    Document.deleted_at.is_not(None)
                    & (Document.deleted_at <= soft_delete_cutoff)
                )
            )
        )
        documents = list(rows.scalars().all())

        deleted_files = 0
        deleted_rows = 0
        for document in documents:
            try:
                if _delete_storage_path(document.storage_path):
                    deleted_files += 1
            except Exception:
                logger.exception(
                    "Failed to delete storage for expired document %s", document.id
                )

            await db.execute(delete(AIJob).where(AIJob.document_id == document.id))
            await db.delete(document)
            deleted_rows += 1

        await db.commit()
        return {
            "status": "completed",
            "deleted_documents": deleted_rows,
            "deleted_files": deleted_files,
            "ran_at": now.isoformat(),
        }


def process_document_analysis(job_id: str | UUID) -> dict[str, Any]:
    """RQ entry point for document analysis."""
    logger.info("Starting document analysis job %s", job_id)
    return asyncio.run(_process_document_analysis_async(job_id))


def batch_document_import(
    org_id: str | UUID,
    user_id: str | UUID,
    documents: list[dict[str, Any] | str],
) -> dict[str, Any]:
    """RQ entry point for batch document import plus immediate analysis."""
    logger.info(
        "Starting batch document import for org=%s user=%s count=%d",
        org_id,
        user_id,
        len(documents),
    )
    return asyncio.run(_batch_document_import_async(org_id, user_id, documents))


def generate_analytics_report(
    org_id: str | UUID,
    report_type: str = "usage",
    days: int = 30,
) -> dict[str, Any]:
    """RQ entry point for organization analytics reports."""
    logger.info(
        "Generating analytics report for org=%s type=%s days=%d",
        org_id,
        report_type,
        days,
    )
    return asyncio.run(
        _generate_analytics_report_async(org_id, report_type=report_type, days=days)
    )


def cleanup_expired_documents() -> dict[str, Any]:
    """Delete expired documents and their backing storage artifacts."""
    logger.info("Running expired document cleanup job")
    return asyncio.run(_cleanup_expired_documents_async())


def send_email_notification(
    to: str | list[str],
    subject: str,
    html: str,
    *,
    reply_to: str | None = None,
    tags: list[dict[str, str]] | None = None,
) -> dict[str, Any]:
    """Worker-safe email wrapper."""
    return send_email(to=to, subject=subject, html=html, reply_to=reply_to, tags=tags)


def send_deletion_request_email_job(to: str, request_id: str) -> dict[str, Any]:
    """Worker entry point for GDPR deletion confirmation email."""
    return send_deletion_request_email(to=to, request_id=request_id)


def send_data_export_email_job(to: str, request_id: str) -> dict[str, Any]:
    """Worker entry point for GDPR export confirmation email."""
    return send_data_export_request_email(to=to, request_id=request_id)


def send_webhook_notification(
    webhook_url: str,
    event_type: str,
    payload: dict[str, Any],
    signing_secret: str | None = None,
) -> dict[str, Any]:
    """Deliver a signed webhook to an external endpoint with SSRF protection."""
    if not _validate_webhook_url(webhook_url):
        return {
            "status": "failed",
            "error": "Webhook delivery blocked: target URL is not allowed",
        }

    body = {
        "event_type": event_type,
        "sent_at": _utcnow().isoformat(),
        "payload": payload,
    }
    headers = {
        "Content-Type": "application/json",
        "User-Agent": "Cogent-Webhook/1.0",
    }

    if signing_secret:
        serialized = json.dumps(body, separators=(",", ":"), sort_keys=True).encode()
        signature = hmac.new(
            signing_secret.encode(), serialized, hashlib.sha256
        ).hexdigest()
        headers["X-Cogent-Signature"] = f"sha256={signature}"

    try:
        response = httpx.post(webhook_url, json=body, headers=headers, timeout=10.0)
        return {
            "status": "success",
            "status_code": response.status_code,
            "event_type": event_type,
        }
    except Exception as exc:
        logger.exception("Webhook delivery failed to %s", webhook_url)
        return {"status": "failed", "error": str(exc), "event_type": event_type}


def run_feedback_retraining(
    *, lookback_days: int = 30, write_snapshot: bool = True
) -> dict[str, Any]:
    """Compatibility wrapper for the feedback retraining service."""
    return run_feedback_retraining_job(
        lookback_days=lookback_days,
        write_snapshot=write_snapshot,
    )
